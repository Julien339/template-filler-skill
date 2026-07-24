"""
Opens the ORIGINAL .pptx (the same file extract_pptx.py read) and writes new
text into specific runs identified by ID in changes.json, then saves the
result to <output.pptx>. Everything not named in changes.json -- shapes,
paragraphs, other runs, all formatting -- is left completely untouched:
writes set run.text on the EXISTING run object in place (python-pptx mutates
the run's own <a:t> text node; its <a:rPr> formatting is never touched),
never delete/recreate a shape or paragraph, and are addressed by ID only --
never by a global search-and-replace across the file, which is exactly how
you'd silently rewrite the wrong occurrence of a repeated value like "12".

changes.json shape: [{"id": "2/5/0/0", "new_text": "Stellantis"}, ...]

An ID with no trailing run index (e.g. "2/5/0" instead of "2/5/0/0") is a
PARAGRAPH-level replacement: the new text goes into the paragraph's first
run and every other run in that paragraph is blanked. Use this when a
paragraph's visible text is split across more XML runs than are meaningful
(a common PowerPoint artifact from autocorrect or a prior manual edit) --
see the content map's "paragraph_text" field for each run's ID.

Usage:
    python apply_pptx.py <template.pptx> <changes.json> <output.pptx>
"""
import sys
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def find_shape_by_path(shapes, path_parts):
    target_id = int(path_parts[0])
    for shape in shapes:
        if shape.shape_id != target_id:
            continue
        if len(path_parts) == 1:
            return shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            return find_shape_by_path(shape.shapes, path_parts[1:])
        return None
    return None


def get_text_frame(shape, rest):
    if rest and rest[0].startswith("r") and "c" in rest[0]:
        row_s, col_s = rest[0][1:].split("c")
        cell = shape.table.cell(int(row_s), int(col_s))
        return cell.text_frame, rest[1:]
    return shape.text_frame, rest


def replace_paragraph_text(paragraph, new_text):
    runs = paragraph.runs
    if not runs:
        paragraph.add_run().text = new_text
        return
    runs[0].text = new_text
    for extra in runs[1:]:
        extra.text = ""


def set_run_text(paragraph, run_idx, new_text):
    runs = paragraph.runs
    if not runs:
        paragraph.add_run().text = new_text
    elif run_idx >= len(runs):
        runs[-1].text = new_text
    else:
        runs[run_idx].text = new_text


def apply_change(prs, change_id, new_text):
    slide_idx_s, shape_path_s, *rest = change_id.split("/")
    slide = prs.slides[int(slide_idx_s)]
    shape = find_shape_by_path(slide.shapes, shape_path_s.split("."))
    if shape is None:
        raise ValueError(f"no shape found for id {change_id!r}")

    tf, rest = get_text_frame(shape, rest)
    para_idx = int(rest[0])
    paragraph = tf.paragraphs[para_idx]

    if len(rest) == 1:
        replace_paragraph_text(paragraph, new_text)
    else:
        set_run_text(paragraph, int(rest[1]), new_text)


def main():
    in_path, changes_path, out_path = sys.argv[1], sys.argv[2], sys.argv[3]
    with open(changes_path, encoding="utf-8") as f:
        changes = json.load(f)

    prs = Presentation(in_path)
    applied, errors = 0, []
    for change in changes:
        try:
            apply_change(prs, change["id"], change["new_text"])
            applied += 1
        except Exception as e:
            errors.append(f"{change['id']!r}: {e}")

    prs.save(out_path)

    print(f"Applied {applied}/{len(changes)} change(s) -> {out_path}")
    if errors:
        print(f"{len(errors)} change(s) failed:", file=sys.stderr)
        for e in errors:
            print(f"  - {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
