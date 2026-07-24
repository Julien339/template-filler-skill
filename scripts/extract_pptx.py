"""
Walks an EXISTING .pptx and emits a flat, addressable content map of every
text run's current text, keyed by a stable structural ID -- the artifact an
agent reads to decide what to change before writing changes.json.

Unlike build_pptx.py (which always creates a brand-new Presentation()), this
opens a real, already-styled file and only READS it -- shape/run access uses
the same python-pptx vocabulary build_pptx.py uses to write
(shape.text_frame.paragraphs[i].runs[j], table.cell(r,c)), just in the read
direction.

ID scheme:
  text frame run:   "{slide}/{shape_id}/{para}/{run}"
  table cell run:   "{slide}/{shape_id}/r{row}c{col}/{para}/{run}"
Group shapes prefix the ID with the group's own shape_id (dot-joined) so
nested shapes stay addressable: "{slide}/{group_id}.{shape_id}/...".
An ID with no trailing run index ("{slide}/{shape_id}/{para}") addresses the
whole paragraph -- apply_pptx.py treats that as "replace the paragraph's
full visible text", which is the escape hatch for paragraphs PowerPoint has
split into more runs (autocorrect, a prior manual edit) than are meaningful.

Merged table cells: table.cell(r,c) returns the same underlying XML element
for every coordinate a merge spans -- detected by element identity so only
the merge-origin cell is emitted once.

python-pptx's paragraph.runs only ever returns <a:r> elements, never <a:fld>
(slide-number/date auto-fields) -- those are simply invisible to this
extractor, which is the safe behavior (nothing to accidentally overwrite).

Picture/media shapes carry no addressable text and are skipped -- image
swapping is out of scope for this version (see README known limits).

Usage:
    python extract_pptx.py <template.pptx> <content_map.json>
"""
import sys
import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def table_origin_cells(table):
    """Yield (row, col, cell) once per distinct cell, skipping the
    duplicate cell references python-pptx returns for merged regions.

    Dedup by the underlying lxml element itself (kept alive in `seen`), not
    by id(cell._tc): python-pptx builds a fresh, unreferenced _Cell/_tc
    wrapper on every table.cell() call, so an id()-based set with no strong
    reference to the wrapper sees CPython immediately recycle that memory
    address for the next wrapper -- producing spurious "duplicate" hits
    between structurally unrelated cells that just happened to reuse the
    same freed address."""
    seen = set()
    for ri in range(len(table.rows)):
        for ci in range(len(table.columns)):
            cell = table.cell(ri, ci)
            if cell._tc in seen:
                continue
            seen.add(cell._tc)
            yield ri, ci, cell


def walk_text_frame(tf, prefix, out_runs):
    for pi, para in enumerate(tf.paragraphs):
        runs = para.runs
        para_text = "".join(r.text for r in runs)
        if not runs:
            out_runs.append({"id": f"{prefix}/{pi}/0", "text": "", "paragraph_text": "", "empty": True})
            continue
        for ri, run in enumerate(runs):
            out_runs.append({"id": f"{prefix}/{pi}/{ri}", "text": run.text, "paragraph_text": para_text})


def walk_table(table, prefix, out_runs):
    for ri, ci, cell in table_origin_cells(table):
        walk_text_frame(cell.text_frame, f"{prefix}/r{ri}c{ci}", out_runs)


def walk_shapes(shapes, slide_idx, group_path, out_runs, out_shapes_meta):
    for shape in shapes:
        shape_path = f"{group_path}.{shape.shape_id}" if group_path else str(shape.shape_id)
        full_prefix = f"{slide_idx}/{shape_path}"

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            out_shapes_meta.append({"id": full_prefix, "name": shape.name, "type": "group"})
            walk_shapes(shape.shapes, slide_idx, shape_path, out_runs, out_shapes_meta)
            continue

        if getattr(shape, "has_table", False):
            out_shapes_meta.append({"id": full_prefix, "name": shape.name, "type": "table"})
            walk_table(shape.table, full_prefix, out_runs)
            continue

        if getattr(shape, "has_text_frame", False):
            out_shapes_meta.append({"id": full_prefix, "name": shape.name, "type": "text"})
            walk_text_frame(shape.text_frame, full_prefix, out_runs)
            continue
        # Pictures and other non-text shapes: not addressable in V1.


def main():
    in_path, out_path = sys.argv[1], sys.argv[2]
    prs = Presentation(in_path)

    slides_out = []
    for si, slide in enumerate(prs.slides):
        runs, shapes_meta = [], []
        walk_shapes(slide.shapes, si, "", runs, shapes_meta)
        slides_out.append({"index": si, "shapes": shapes_meta, "runs": runs})

    content_map = {
        "format": "pptx",
        "source": in_path,
        "slide_count": len(prs.slides),
        "slides": slides_out,
    }
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(content_map, f, ensure_ascii=False, indent=2)

    total_runs = sum(len(s["runs"]) for s in slides_out)
    print(f"Extracted {len(slides_out)} slide(s), {total_runs} addressable run(s) -> {out_path}")


if __name__ == "__main__":
    main()
