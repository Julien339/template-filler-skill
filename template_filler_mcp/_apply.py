"""Text application for PPTX and DOCX files.

Opens the ORIGINAL template file and writes new text into specific runs
identified by ID, then saves the result. Everything not named in the changes
list — shapes, paragraphs, other runs, all formatting — is left untouched.
"""

from __future__ import annotations

from typing import Any

from docx import Document
from docx.table import Table
from docx.text.paragraph import Paragraph
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ._path_validation import validate_input_file, validate_output_file

from ._path_validation import validate_input_file, validate_output_file

# ── shared helpers ─────────────────────────────────────────────────────


def _replace_paragraph_text(paragraph, new_text):
    """Replace the full visible text of a paragraph.

    The new text goes into the first run; every other run is blanked.
    This is the escape hatch for paragraphs PowerPoint/Word have split into
    more XML runs than are meaningful (autocorrect, prior manual edits).
    """
    runs = paragraph.runs
    if not runs:
        paragraph.add_run(new_text)
        return
    runs[0].text = new_text
    for extra in runs[1:]:
        extra.text = ""


def _set_run_text(paragraph, run_idx, new_text):
    """Set the text of a single run by index.

    If the paragraph has no runs, a new one is created. If run_idx is past
    the last run, the last run is used (best-effort fallback for edge cases).
    """
    runs = paragraph.runs
    if not runs:
        paragraph.add_run(new_text)
    elif run_idx >= len(runs):
        runs[-1].text = new_text
    else:
        runs[run_idx].text = new_text


# ── PPTX application ───────────────────────────────────────────────────


def _find_shape_by_path(shapes, path_parts):
    target_id = int(path_parts[0])
    for shape in shapes:
        if shape.shape_id != target_id:
            continue
        if len(path_parts) == 1:
            return shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            return _find_shape_by_path(shape.shapes, path_parts[1:])
        return None
    return None


def _get_text_frame(shape, rest):
    if rest and rest[0].startswith("r") and "c" in rest[0]:
        row_s, col_s = rest[0][1:].split("c")
        cell = shape.table.cell(int(row_s), int(col_s))
        return cell.text_frame, rest[1:]
    return shape.text_frame, rest


def _apply_pptx_change(prs, change_id, new_text):
    slide_idx_s, shape_path_s, *rest = change_id.split("/")
    slide = prs.slides[int(slide_idx_s)]
    shape = _find_shape_by_path(slide.shapes, shape_path_s.split("."))
    if shape is None:
        raise ValueError(f"no shape found for id {change_id!r}")

    tf, rest = _get_text_frame(shape, rest)
    para_idx = int(rest[0])
    paragraph = tf.paragraphs[para_idx]

    if len(rest) == 1:
        _replace_paragraph_text(paragraph, new_text)
    else:
        _set_run_text(paragraph, int(rest[1]), new_text)


def apply_pptx_changes(template_path: str, changes: list[dict[str, str]], output_path: str) -> dict[str, Any]:
    """Apply text changes to a PPTX template.

    Args:
        template_path: Path to the original .pptx file.
        changes: List of {"id": "...", "new_text": "..."} objects.
            ID format: {slide}/{shape_id}/{paragraph}/{run} for runs,
            or {slide}/{shape_id}/{paragraph} for whole paragraphs.
        output_path: Path to save the modified file.

    Returns:
        {"applied": N, "total": M, "failed": [...], "output_path": "..."}
    """
    # Validate paths for security
    validated_template = validate_input_file(template_path, allowed_extensions=('.pptx',))
    validated_output = validate_output_file(output_path, allowed_extensions=('.pptx',))

    prs = Presentation(str(validated_template))
    applied, errors = 0, []
    for change in changes:
        try:
            _apply_pptx_change(prs, change["id"], change["new_text"])
            applied += 1
        except Exception as e:
            errors.append({"id": change["id"], "error": str(e)})

    prs.save(str(validated_output))

    result = {
        "applied": applied,
        "total": len(changes),
        "failed": errors,
        "output_path": str(validated_output),
    }
    return result


# ── DOCX application ───────────────────────────────────────────────────


def _resolve_paragraph(blocks, doc, path_parts):
    """Resolve a change id to (paragraph, remaining_path_parts).

    path_parts examples:
        ["3"], ["3","0"]               — body paragraph runs
        ["5","r1c2","0"], ["5","r1c2","0","0"] — table cell runs
        ["h0","0"], ["h0","0","0"]     — header runs
        ["f0","0"], ["f0","0","0"]     — footer runs
    """
    head = path_parts[0]

    if head[0] in ("h", "f"):
        section_idx = int(head[1:])
        section = doc.sections[section_idx]
        part = section.header if head[0] == "h" else section.footer
        para_idx = int(path_parts[1])
        return part.paragraphs[para_idx], path_parts[2:]

    block = blocks[int(head)]
    rest = path_parts[1:]

    if isinstance(block, Table):
        row_s, col_s = rest[0][1:].split("c")
        cell = block.cell(int(row_s), int(col_s))
        para_idx = int(rest[1])
        return cell.paragraphs[para_idx], rest[2:]

    assert isinstance(block, Paragraph)
    return block, rest


def _apply_docx_change(blocks, doc, change_id, new_text):
    paragraph, rest = _resolve_paragraph(blocks, doc, change_id.split("/"))
    if len(rest) == 0:
        _replace_paragraph_text(paragraph, new_text)
    else:
        _set_run_text(paragraph, int(rest[0]), new_text)


def apply_docx_changes(template_path: str, changes: list[dict[str, str]], output_path: str) -> dict[str, Any]:
    """Apply text changes to a DOCX template.

    Args:
        template_path: Path to the original .docx file.
        changes: List of {"id": "...", "new_text": "..."} objects.
        output_path: Path to save the modified file.

    Returns:
        {"applied": N, "total": M, "failed": [...], "output_path": "..."}
    """
    # Validate paths for security
    validated_template = validate_input_file(template_path, allowed_extensions=('.docx',))
    validated_output = validate_output_file(output_path, allowed_extensions=('.docx',))

    doc = Document(str(validated_template))
    blocks = list(doc.iter_inner_content())

    applied, errors = 0, []
    for change in changes:
        try:
            _apply_docx_change(blocks, doc, change["id"], change["new_text"])
            applied += 1
        except Exception as e:
            errors.append({"id": change["id"], "error": str(e)})

    doc.save(str(validated_output))

    result = {
        "applied": applied,
        "total": len(changes),
        "failed": errors,
        "output_path": str(validated_output),
    }
    return result
