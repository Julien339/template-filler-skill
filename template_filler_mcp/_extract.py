"""Content extraction for PPTX and DOCX files.

Extracts all text runs from a template file into a flat, addressable JSON
structure keyed by stable structural IDs — the artifact an agent reads to
decide what to change before constructing a changes list.
"""

from __future__ import annotations

from typing import Any

from docx import Document
from docx.oxml.ns import qn
from docx.table import Table
from docx.text.paragraph import Paragraph
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ── shared: merged-table-cell deduplication ────────────────────────────


def _table_origin_cells(table):
    """Yield (row, col, cell) once per distinct cell, skipping duplicate
    references python-pptx/python-docx return for merged regions.

    Dedup by the underlying lxml element itself (kept alive in ``seen``), not
    by id(cell._tc): python-pptx builds a fresh, unreferenced wrapper on every
    ``table.cell()`` call, so an id()-based set with no strong reference to
    the wrapper sees CPython immediately recycle that memory address for the
    next wrapper — producing spurious "duplicate" hits between structurally
    unrelated cells that just happened to reuse the same freed address.
    """
    seen = set()
    for ri in range(len(table.rows)):
        for ci in range(len(table.columns)):
            cell = table.cell(ri, ci)
            if cell._tc in seen:
                continue
            seen.add(cell._tc)
            yield ri, ci, cell


# ── PPTX extraction ────────────────────────────────────────────────────


def _walk_text_frame(tf, prefix, out_runs):
    for pi, para in enumerate(tf.paragraphs):
        runs = para.runs
        para_text = "".join(r.text for r in runs)
        if not runs:
            out_runs.append({"id": f"{prefix}/{pi}/0", "text": "", "paragraph_text": "", "empty": True})
            continue
        for ri, run in enumerate(runs):
            out_runs.append({"id": f"{prefix}/{pi}/{ri}", "text": run.text, "paragraph_text": para_text})


def _walk_table(table, prefix, out_runs):
    for ri, ci, cell in _table_origin_cells(table):
        _walk_text_frame(cell.text_frame, f"{prefix}/r{ri}c{ci}", out_runs)


def _walk_shapes(shapes, slide_idx, group_path, out_runs, out_shapes_meta):
    for shape in shapes:
        shape_path = f"{group_path}.{shape.shape_id}" if group_path else str(shape.shape_id)
        full_prefix = f"{slide_idx}/{shape_path}"

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            out_shapes_meta.append({"id": full_prefix, "name": shape.name, "type": "group"})
            _walk_shapes(shape.shapes, slide_idx, shape_path, out_runs, out_shapes_meta)
            continue

        if getattr(shape, "has_table", False):
            out_shapes_meta.append({"id": full_prefix, "name": shape.name, "type": "table"})
            _walk_table(shape.table, full_prefix, out_runs)
            continue

        if getattr(shape, "has_text_frame", False):
            out_shapes_meta.append({"id": full_prefix, "name": shape.name, "type": "text"})
            _walk_text_frame(shape.text_frame, full_prefix, out_runs)
            continue
        # Pictures and other non-text shapes: not addressable in V1.


def extract_pptx_content(template_path: str) -> dict[str, Any]:
    """Extract all text runs from a PPTX file.

    Returns a content map with:
      - format: "pptx"
      - source: template_path
      - slide_count: int
      - slides: list of {index, shapes, runs}

    Each run has: id, text, paragraph_text, and optionally empty=True.
    ID scheme: {slide}/{shape_id}/{paragraph}/{run}
    Table cells: {slide}/{shape_id}/r{row}c{col}/{paragraph}/{run}
    Group shapes prefix the group shape_id (dot-joined).

    Merged table cells are deduplicated — only the merge-origin cell is emitted.
    """
    prs = Presentation(template_path)

    slides_out = []
    for si, slide in enumerate(prs.slides):
        runs, shapes_meta = [], []
        _walk_shapes(slide.shapes, si, "", runs, shapes_meta)
        slides_out.append({"index": si, "shapes": shapes_meta, "runs": runs})

    content_map = {
        "format": "pptx",
        "source": template_path,
        "slide_count": len(prs.slides),
        "slides": slides_out,
    }
    return content_map


# ── DOCX extraction ────────────────────────────────────────────────────

W = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"


def _paragraph_field_flags(paragraph):
    """One bool per paragraph.runs entry: True if that run is part of a
    Word field (date/page-number/TOC) rather than plain authored text."""
    flags = []
    in_field = False
    for run in paragraph.runs:
        r = run._r
        fld_char = r.find(qn("w:fldChar"))
        instr = r.find(qn("w:instrText"))
        if fld_char is not None:
            fld_type = fld_char.get(qn("w:fldCharType"))
            flags.append(True)
            in_field = fld_type != "end"
            continue
        if instr is not None:
            flags.append(True)
            continue
        flags.append(in_field)
    return flags


def _emit_paragraph_runs(paragraph, prefix, out_runs):
    runs = paragraph.runs
    if not runs:
        out_runs.append({"id": f"{prefix}/0", "text": "", "paragraph_text": "", "empty": True})
        return
    field_flags = _paragraph_field_flags(paragraph)
    para_text = "".join(r.text for r in runs)
    for ri, (run, is_field) in enumerate(zip(runs, field_flags)):
        entry = {"id": f"{prefix}/{ri}", "text": run.text, "paragraph_text": para_text}
        if is_field:
            entry["field"] = True
        out_runs.append(entry)


def _emit_table(table, prefix, out_runs):
    for ri, ci, cell in _table_origin_cells(table):
        for pi, para in enumerate(cell.paragraphs):
            _emit_paragraph_runs(para, f"{prefix}/r{ri}c{ci}/{pi}", out_runs)


def _emit_header_footer(part, label, out_runs):
    if part is None or part.is_linked_to_previous:
        return
    for pi, para in enumerate(part.paragraphs):
        _emit_paragraph_runs(para, f"{label}/{pi}", out_runs)


def extract_docx_content(template_path: str) -> dict[str, Any]:
    """Extract all text runs from a DOCX file.

    Returns a content map with:
      - format: "docx"
      - source: template_path
      - section_count: int
      - body_runs: list of runs from body paragraphs and tables
      - header_footer_runs: list of runs from default headers/footers

    Each run has: id, text, paragraph_text, and optionally field=True, empty=True.
    ID scheme for body: {block}/{run} (paragraph), {block}/r{row}c{col}/{para}/{run} (table)
    ID scheme for headers/footers: h{section}/{para}/{run}, f{section}/{para}/{run}

    Merged table cells are deduplicated. Field-coded runs (auto-fields like
    dates, page numbers, TOC entries) are flagged as read-only.
    """
    doc = Document(template_path)

    body_runs = []
    for bi, block in enumerate(doc.iter_inner_content()):
        if isinstance(block, Paragraph):
            _emit_paragraph_runs(block, str(bi), body_runs)
        elif isinstance(block, Table):
            _emit_table(block, str(bi), body_runs)

    header_footer_runs = []
    for si, section in enumerate(doc.sections):
        _emit_header_footer(section.header, f"h{si}", header_footer_runs)
        _emit_header_footer(section.footer, f"f{si}", header_footer_runs)

    content_map = {
        "format": "docx",
        "source": template_path,
        "section_count": len(doc.sections),
        "body_runs": body_runs,
        "header_footer_runs": header_footer_runs,
    }
    return content_map
